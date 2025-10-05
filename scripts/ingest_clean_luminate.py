#!/usr/bin/env python3
"""
Luminate AI - Blackboard LMS Data Ingestion & Cleaning Pipeline
================================================================

This script processes Blackboard course exports and transforms them into a structured
dataset suitable for ChromaDB + LangGraph ingestion with live-link mapping to 
Blackboard URLs.

Author: Luminate AI Data Engineering Team
Date: October 4, 2025
"""

import json
import logging
import re
import xml.etree.ElementTree as ET
from collections import defaultdict
from dataclasses import dataclass, field, asdict
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Set
from datetime import datetime
import hashlib

# Third-party imports
try:
    from bs4 import BeautifulSoup
    from pypdf import PdfReader
    from docx import Document as DocxDocument
    from pptx import Presentation
    import chardet
    from tqdm import tqdm
except ImportError as e:
    print(f"Missing required library: {e}")
    print("Install with: pip install beautifulsoup4 pypdf python-docx python-pptx chardet tqdm")
    exit(1)


# ============================================================================
# CONFIGURATION
# ============================================================================

@dataclass
class Config:
    """Central configuration for the ingestion pipeline."""
    
    # Course identifiers
    course_id: str = "_29430_1"
    course_name: str = "Luminate"
    base_url: str = "https://luminate.centennialcollege.ca/ultra/courses"
    
    # Paths
    source_dir: Path = Path("extracted/ExportFile_COMP237_INP")
    output_dir: Path = Path("cleaned")
    graph_dir: Path = Path("graph_seed")
    log_dir: Path = Path("logs")
    
    # Text processing
    chunk_min_tokens: int = 500
    chunk_max_tokens: int = 800
    chunk_overlap_percent: float = 0.5
    chars_per_token: float = 4.0  # Rough approximation
    
    # File processing
    supported_extensions: Set[str] = field(default_factory=lambda: {
        '.html', '.htm', '.pdf', '.docx', '.pptx', '.txt', '.xml', '.md', '.dat',
        '.csv', '.py'  # Added CSV and Python files from the export
    })
    
    # Blackboard patterns
    bb_id_pattern: re.Pattern = field(default_factory=lambda: re.compile(r'_(\d+)_1'))
    
    # Boilerplate removal patterns
    boilerplate_patterns: List[str] = field(default_factory=lambda: [
        r'Blackboard\s+Learn',
        r'Breadcrumb',
        r'Last\s+edited',
        r'Print\s+this\s+page',
        r'Show\s+all\s+content',
        r'Hide\s+all\s+content',
        r'View\s+course\s+map',
        r'ultraDocumentBody',
        r'@X@EmbeddedFile\.requestUrlStub@X@',
    ])
    
    def __post_init__(self):
        """Convert string patterns to compiled regexes."""
        self.boilerplate_patterns = [
            re.compile(p, re.IGNORECASE) for p in self.boilerplate_patterns
        ]


# ============================================================================
# DATA MODELS
# ============================================================================

@dataclass
class Chunk:
    """Represents a text chunk with metadata."""
    chunk_id: str
    content: str
    tags: List[str]
    live_lms_url: Optional[str] = None
    token_count: int = 0
    chunk_index: int = 0
    total_chunks: int = 0


@dataclass
class DocumentMetadata:
    """Metadata for a processed document."""
    course_id: str
    course_name: str
    file_name: str
    content_type: str
    original_file_path: str
    module_name: Optional[str] = None
    bb_doc_id: Optional[str] = None
    live_lms_url: Optional[str] = None
    title: Optional[str] = None
    created_date: Optional[str] = None
    updated_date: Optional[str] = None
    parent_id: Optional[str] = None


@dataclass
class ProcessedDocument:
    """Complete processed document with chunks and metadata."""
    metadata: DocumentMetadata
    chunks: List[Chunk]
    raw_text_length: int = 0
    total_tokens: int = 0


@dataclass
class GraphLink:
    """Represents a relationship between documents."""
    source: str
    target: str
    relation: str
    metadata: Dict = field(default_factory=dict)


# ============================================================================
# LOGGING SETUP
# ============================================================================

class IngestionLogger:
    """Manages logging for the ingestion pipeline."""
    
    def __init__(self, log_dir: Path):
        self.log_dir = log_dir
        self.log_dir.mkdir(parents=True, exist_ok=True)
        
        # Setup main logger
        self.logger = logging.getLogger('LuminateIngestion')
        self.logger.setLevel(logging.DEBUG)
        
        # File handler for all logs
        fh = logging.FileHandler(log_dir / 'ingestion.log', mode='w', encoding='utf-8')
        fh.setLevel(logging.DEBUG)
        
        # Console handler for INFO and above
        ch = logging.StreamHandler()
        ch.setLevel(logging.INFO)
        
        # Formatter
        formatter = logging.Formatter(
            '%(asctime)s - %(name)s - %(levelname)s - %(message)s',
            datefmt='%Y-%m-%d %H:%M:%S'
        )
        fh.setFormatter(formatter)
        ch.setFormatter(formatter)
        
        self.logger.addHandler(fh)
        self.logger.addHandler(ch)
        
        # Issue tracking file
        self.issues_file = log_dir / 'ingest_issues.txt'
        self.issues = []
    
    def log_issue(self, severity: str, file_path: str, message: str):
        """Log an issue to the issues file."""
        issue = f"[{severity}] {file_path}: {message}"
        self.issues.append(issue)
        self.logger.warning(issue)
    
    def save_issues(self):
        """Save all issues to file."""
        with open(self.issues_file, 'w', encoding='utf-8') as f:
            f.write(f"Ingestion Issues Report\n")
            f.write(f"Generated: {datetime.now().isoformat()}\n")
            f.write(f"{'=' * 80}\n\n")
            for issue in self.issues:
                f.write(f"{issue}\n")


# ============================================================================
# FILE PARSERS
# ============================================================================

class FileParser:
    """Base class for file parsing."""
    
    def __init__(self, config: Config, logger: IngestionLogger):
        self.config = config
        self.logger = logger
    
    def clean_text(self, text: str) -> str:
        """Clean and normalize text."""
        if not text:
            return ""
        
        # Remove boilerplate patterns
        for pattern in self.config.boilerplate_patterns:
            text = pattern.sub('', text)
        
        # Normalize whitespace
        text = re.sub(r'\s+', ' ', text)
        text = re.sub(r'\n\s*\n\s*\n+', '\n\n', text)
        
        return text.strip()
    
    def extract_bb_id(self, text: str, filename: str = "") -> Optional[str]:
        """Extract Blackboard document ID from text or filename."""
        # Try filename first
        if filename:
            match = self.config.bb_id_pattern.search(filename)
            if match:
                return f"_{match.group(1)}_1"
        
        # Try text content
        match = self.config.bb_id_pattern.search(text)
        if match:
            return f"_{match.group(1)}_1"
        
        return None
    
    def construct_live_url(self, bb_doc_id: str) -> str:
        """Construct live LMS URL from Blackboard document ID."""
        return (
            f"{self.config.base_url}/{self.config.course_id}/outline/edit/"
            f"document/{bb_doc_id}?courseId={self.config.course_id}&view=content&state=view"
        )


class DatFileParser(FileParser):
    """Parser for Blackboard .dat files (XML format)."""
    
    def parse(self, file_path: Path) -> Tuple[str, Dict]:
        """Parse a .dat file and extract content and metadata."""
        try:
            tree = ET.parse(file_path)
            root = tree.getroot()
            
            metadata = {}
            content_parts = []
            
            # Extract metadata
            if 'id' in root.attrib:
                metadata['bb_doc_id'] = root.attrib['id']
            
            # Extract title
            title_elem = root.find('.//TITLE')
            if title_elem is not None and 'value' in title_elem.attrib:
                metadata['title'] = title_elem.attrib['value']
            
            # Extract dates
            created_elem = root.find('.//CREATED')
            if created_elem is not None and 'value' in created_elem.attrib:
                metadata['created_date'] = created_elem.attrib['value']
            
            updated_elem = root.find('.//UPDATED')
            if updated_elem is not None and 'value' in updated_elem.attrib:
                metadata['updated_date'] = updated_elem.attrib['value']
            
            # Extract parent ID for hierarchy
            parent_elem = root.find('.//PARENTID')
            if parent_elem is not None and 'value' in parent_elem.attrib:
                metadata['parent_id'] = parent_elem.attrib['value']
            
            # Extract HTML content from BODY/TEXT
            text_elem = root.find('.//BODY/TEXT')
            if text_elem is not None and text_elem.text:
                # Parse HTML content
                html_content = text_elem.text
                soup = BeautifulSoup(html_content, 'html.parser')
                
                # Extract text with basic markdown conversion
                for heading in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6']):
                    level = int(heading.name[1])
                    heading_text = heading.get_text(strip=True)
                    content_parts.append(f"\n{'#' * level} {heading_text}\n")
                
                # Get remaining text
                text_content = soup.get_text(separator='\n', strip=True)
                content_parts.append(text_content)
            
            # Extract description
            desc_elem = root.find('.//DESCRIPTION')
            if desc_elem is not None and 'value' in desc_elem.attrib:
                desc = desc_elem.attrib['value']
                if desc:
                    content_parts.append(desc)
            
            full_text = '\n\n'.join(content_parts)
            return self.clean_text(full_text), metadata
            
        except ET.ParseError as e:
            self.logger.log_issue('ERROR', str(file_path), f'XML parse error: {e}')
            return "", {}
        except Exception as e:
            self.logger.log_issue('ERROR', str(file_path), f'Unexpected error: {e}')
            return "", {}


class HtmlParser(FileParser):
    """Parser for HTML files."""
    
    def parse(self, file_path: Path) -> str:
        """Parse HTML file and extract clean text."""
        try:
            # Detect encoding
            with open(file_path, 'rb') as f:
                raw_data = f.read()
                detected = chardet.detect(raw_data)
                encoding = detected['encoding'] or 'utf-8'
            
            # Read with detected encoding
            with open(file_path, 'r', encoding=encoding, errors='ignore') as f:
                html_content = f.read()
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Remove script and style elements
            for element in soup(['script', 'style', 'nav', 'footer', 'header']):
                element.decompose()
            
            # Convert headings to markdown
            text_parts = []
            for heading in soup.find_all(['h1', 'h2', 'h3']):
                level = int(heading.name[1])
                heading_text = heading.get_text(strip=True)
                text_parts.append(f"\n{'#' * level} {heading_text}\n")
            
            # Extract main text
            text = soup.get_text(separator='\n', strip=True)
            
            return self.clean_text(text)
            
        except Exception as e:
            self.logger.log_issue('ERROR', str(file_path), f'HTML parse error: {e}')
            return ""


class PdfParser(FileParser):
    """Parser for PDF files."""
    
    def parse(self, file_path: Path) -> str:
        """Parse PDF file and extract text."""
        try:
            reader = PdfReader(file_path)
            text_parts = []
            
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
            
            full_text = '\n\n'.join(text_parts)
            return self.clean_text(full_text)
            
        except Exception as e:
            self.logger.log_issue('ERROR', str(file_path), f'PDF parse error: {e}')
            return ""


class DocxParser(FileParser):
    """Parser for DOCX files."""
    
    def parse(self, file_path: Path) -> str:
        """Parse DOCX file and extract text."""
        try:
            doc = DocxDocument(file_path)
            text_parts = []
            
            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if text:
                    # Check if it's a heading
                    if paragraph.style.name.startswith('Heading'):
                        level = int(paragraph.style.name.split()[-1]) if paragraph.style.name.split()[-1].isdigit() else 1
                        text_parts.append(f"\n{'#' * min(level, 3)} {text}\n")
                    else:
                        text_parts.append(text)
            
            full_text = '\n\n'.join(text_parts)
            return self.clean_text(full_text)
            
        except Exception as e:
            self.logger.log_issue('ERROR', str(file_path), f'DOCX parse error: {e}')
            return ""


class PptxParser(FileParser):
    """Parser for PPTX files."""
    
    def parse(self, file_path: Path) -> str:
        """Parse PPTX file and extract text."""
        try:
            prs = Presentation(file_path)
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text_parts.append(f"\n## Slide {slide_num}\n")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_parts.append(shape.text)
            
            full_text = '\n\n'.join(text_parts)
            return self.clean_text(full_text)
            
        except Exception as e:
            self.logger.log_issue('ERROR', str(file_path), f'PPTX parse error: {e}')
            return ""


class TextParser(FileParser):
    """Parser for plain text files."""
    
    def parse(self, file_path: Path) -> str:
        """Parse text file with encoding detection."""
        try:
            # Detect encoding
            with open(file_path, 'rb') as f:
                raw_data = f.read()
                detected = chardet.detect(raw_data)
                encoding = detected['encoding'] or 'utf-8'
            
            # Read with detected encoding
            with open(file_path, 'r', encoding=encoding, errors='ignore') as f:
                text = f.read()
            
            return self.clean_text(text)
            
        except Exception as e:
            self.logger.log_issue('ERROR', str(file_path), f'Text parse error: {e}')
            return ""


# ============================================================================
# TEXT CHUNKING
# ============================================================================

class TextChunker:
    """Handles intelligent text chunking with overlap."""
    
    def __init__(self, config: Config):
        self.config = config
    
    def estimate_tokens(self, text: str) -> int:
        """Estimate token count from text."""
        return int(len(text) / self.config.chars_per_token)
    
    def chunk_text(self, text: str, doc_id: str, live_url: Optional[str] = None, 
                   tags: List[str] = None) -> List[Chunk]:
        """Split text into overlapping chunks."""
        if not text:
            return []
        
        if tags is None:
            tags = []
        
        # Calculate chunk sizes in characters
        min_chars = int(self.config.chunk_min_tokens * self.config.chars_per_token)
        max_chars = int(self.config.chunk_max_tokens * self.config.chars_per_token)
        overlap_chars = int(max_chars * self.config.chunk_overlap_percent)
        
        chunks = []
        start = 0
        chunk_index = 0
        
        while start < len(text):
            # Calculate end position
            end = start + max_chars
            
            # If this isn't the last chunk, try to break at a sentence or paragraph
            if end < len(text):
                # Look for paragraph break first
                paragraph_break = text.rfind('\n\n', start + min_chars, end)
                if paragraph_break != -1:
                    end = paragraph_break
                else:
                    # Look for sentence break
                    sentence_break = max(
                        text.rfind('. ', start + min_chars, end),
                        text.rfind('.\n', start + min_chars, end),
                        text.rfind('! ', start + min_chars, end),
                        text.rfind('? ', start + min_chars, end)
                    )
                    if sentence_break != -1:
                        end = sentence_break + 1
            
            # Extract chunk
            chunk_text = text[start:end].strip()
            
            if chunk_text:
                chunk_id = f"{doc_id}_chunk_{chunk_index:03d}"
                token_count = self.estimate_tokens(chunk_text)
                
                chunk = Chunk(
                    chunk_id=chunk_id,
                    content=chunk_text,
                    tags=tags.copy(),
                    live_lms_url=live_url,
                    token_count=token_count,
                    chunk_index=chunk_index,
                    total_chunks=0  # Will be updated later
                )
                chunks.append(chunk)
                chunk_index += 1
            
            # Move start position with overlap
            start = end - overlap_chars
            
            # Ensure we make progress
            if start <= chunks[-1].content.rfind(chunk_text[:50]) if chunks else 0:
                start = end
        
        # Update total_chunks for all chunks
        for chunk in chunks:
            chunk.total_chunks = len(chunks)
        
        return chunks


# ============================================================================
# MANIFEST PARSER
# ============================================================================

class ManifestParser:
    """Parses imsmanifest.xml to build course structure."""
    
    def __init__(self, manifest_path: Path, logger: IngestionLogger):
        self.manifest_path = manifest_path
        self.logger = logger
        self.structure = {}
        self.resource_map = {}
    
    def parse(self) -> Dict:
        """Parse the manifest and build structure map."""
        try:
            tree = ET.parse(self.manifest_path)
            root = tree.getroot()
            
            # Parse resources
            resources = root.find('.//{http://www.imsglobal.org/xsd/imscp_v1p1}resources')
            if resources is None:
                resources = root.find('.//resources')
            
            if resources is not None:
                for resource in resources.findall('.//{http://www.imsglobal.org/xsd/imscp_v1p1}resource'):
                    if resource is None:
                        resource = resources.find('.//resource')
                    if resource is not None:
                        res_id = resource.get('identifier', '')
                        href = resource.get('href', '')
                        if res_id and href:
                            self.resource_map[res_id] = href
            
            # Parse organization structure
            org = root.find('.//{http://www.imsglobal.org/xsd/imscp_v1p1}organization')
            if org is None:
                org = root.find('.//organization')
            
            if org is not None:
                self._parse_item(org, [])
            
            self.logger.logger.info(f"Parsed manifest: {len(self.resource_map)} resources")
            return self.structure
            
        except Exception as e:
            self.logger.log_issue('ERROR', str(self.manifest_path), f'Manifest parse error: {e}')
            return {}
    
    def _parse_item(self, item, path: List[str]):
        """Recursively parse item structure."""
        # Handle both namespaced and non-namespaced elements
        title_elem = item.find('.//{http://www.imsglobal.org/xsd/imscp_v1p1}title')
        if title_elem is None:
            title_elem = item.find('.//title')
        
        title = title_elem.text if title_elem is not None else "Untitled"
        
        identifierref = item.get('identifierref', '')
        
        current_path = path + [title]
        
        if identifierref:
            self.structure[identifierref] = {
                'path': current_path,
                'title': title
            }
        
        # Process children
        for child in item.findall('.//{http://www.imsglobal.org/xsd/imscp_v1p1}item'):
            if child is None:
                child = item.find('.//item')
            if child is not None and child != item:
                self._parse_item(child, current_path)


# ============================================================================
# GRAPH BUILDER
# ============================================================================

class GraphBuilder:
    """Builds relationship graph between documents."""
    
    def __init__(self, config: Config):
        self.config = config
        self.links: List[GraphLink] = []
        self.document_order: Dict[str, List[str]] = defaultdict(list)
    
    def add_document(self, doc_id: str, parent_id: Optional[str], module_name: Optional[str]):
        """Register a document in the graph."""
        if parent_id:
            # Parent-child relationship
            self.links.append(GraphLink(
                source=parent_id,
                target=doc_id,
                relation="CONTAINS",
                metadata={"type": "hierarchy"}
            ))
        
        if module_name:
            # Track document order within modules
            self.document_order[module_name].append(doc_id)
    
    def build_sequential_links(self):
        """Create NEXT_IN_MODULE and PREV_IN_MODULE links."""
        for module, doc_ids in self.document_order.items():
            for i in range(len(doc_ids) - 1):
                self.links.append(GraphLink(
                    source=doc_ids[i],
                    target=doc_ids[i + 1],
                    relation="NEXT_IN_MODULE",
                    metadata={"module": module}
                ))
                self.links.append(GraphLink(
                    source=doc_ids[i + 1],
                    target=doc_ids[i],
                    relation="PREV_IN_MODULE",
                    metadata={"module": module}
                ))
    
    def save(self, output_path: Path):
        """Save graph links to JSON."""
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        links_data = [
            {
                "source": link.source,
                "target": link.target,
                "relation": link.relation,
                "metadata": link.metadata
            }
            for link in self.links
        ]
        
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(links_data, f, indent=2, ensure_ascii=False)


# ============================================================================
# MAIN INGESTION PIPELINE
# ============================================================================

class IngestionPipeline:
    """Main orchestrator for the ingestion pipeline."""
    
    def __init__(self, config: Config):
        self.config = config
        self.logger = IngestionLogger(config.log_dir)
        
        # Initialize parsers
        self.dat_parser = DatFileParser(config, self.logger)
        self.html_parser = HtmlParser(config, self.logger)
        self.pdf_parser = PdfParser(config, self.logger)
        self.docx_parser = DocxParser(config, self.logger)
        self.pptx_parser = PptxParser(config, self.logger)
        self.text_parser = TextParser(config, self.logger)
        
        self.chunker = TextChunker(config)
        self.graph_builder = GraphBuilder(config)
        
        # Statistics
        self.stats = {
            'total_files': 0,
            'processed_files': 0,
            'skipped_files': 0,
            'total_chunks': 0,
            'total_tokens': 0,
            'by_module': defaultdict(lambda: {'files': 0, 'chunks': 0, 'tokens': 0}),
            'by_type': defaultdict(int),
            'errors': 0
        }
    
    def process_file(self, file_path: Path) -> Optional[ProcessedDocument]:
        """Process a single file."""
        suffix = file_path.suffix.lower()
        
        # Parse based on file type
        text = ""
        metadata_dict = {}
        
        if suffix == '.dat':
            text, metadata_dict = self.dat_parser.parse(file_path)
        elif suffix in ['.html', '.htm']:
            text = self.html_parser.parse(file_path)
        elif suffix == '.pdf':
            text = self.pdf_parser.parse(file_path)
        elif suffix == '.docx':
            text = self.docx_parser.parse(file_path)
        elif suffix == '.pptx':
            text = self.pptx_parser.parse(file_path)
        elif suffix in ['.txt', '.md', '.xml']:
            text = self.text_parser.parse(file_path)
        else:
            self.logger.log_issue('SKIP', str(file_path), f'Unsupported file type: {suffix}')
            return None
        
        if not text:
            self.logger.log_issue('SKIP', str(file_path), 'No text content extracted')
            return None
        
        # Determine module name from path
        parts = file_path.relative_to(self.config.source_dir).parts
        module_name = parts[0] if len(parts) > 1 else "Root"
        
        # Extract Blackboard ID
        bb_doc_id = metadata_dict.get('bb_doc_id') or self.dat_parser.extract_bb_id(text, file_path.name)
        
        # Construct live URL
        live_url = None
        if bb_doc_id:
            live_url = self.dat_parser.construct_live_url(bb_doc_id)
        
        # Create metadata
        metadata = DocumentMetadata(
            course_id=self.config.course_id,
            course_name=self.config.course_name,
            file_name=file_path.name,
            content_type=suffix,
            original_file_path=str(file_path.relative_to(self.config.source_dir)),
            module_name=module_name,
            bb_doc_id=bb_doc_id,
            live_lms_url=live_url,
            title=metadata_dict.get('title'),
            created_date=metadata_dict.get('created_date'),
            updated_date=metadata_dict.get('updated_date'),
            parent_id=metadata_dict.get('parent_id')
        )
        
        # Generate unique document ID
        doc_id = bb_doc_id or f"doc_{hashlib.md5(str(file_path).encode()).hexdigest()[:12]}"
        
        # Create chunks
        tags = [module_name]
        if metadata.title:
            tags.append(metadata.title)
        
        chunks = self.chunker.chunk_text(text, doc_id, live_url, tags)
        
        # Create processed document
        doc = ProcessedDocument(
            metadata=metadata,
            chunks=chunks,
            raw_text_length=len(text),
            total_tokens=sum(c.token_count for c in chunks)
        )
        
        # Add to graph
        self.graph_builder.add_document(doc_id, metadata.parent_id, module_name)
        
        return doc
    
    def save_document(self, doc: ProcessedDocument):
        """Save processed document to JSON."""
        # Create output path mirroring source structure
        relative_path = Path(doc.metadata.original_file_path)
        output_path = self.config.output_dir / relative_path.parent / f"{relative_path.stem}.json"
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        # Prepare output data
        output_data = {
            "course_id": doc.metadata.course_id,
            "course_name": doc.metadata.course_name,
            "module": doc.metadata.module_name,
            "file_name": doc.metadata.file_name,
            "content_type": doc.metadata.content_type,
            "bb_doc_id": doc.metadata.bb_doc_id,
            "live_lms_url": doc.metadata.live_lms_url,
            "title": doc.metadata.title,
            "created_date": doc.metadata.created_date,
            "updated_date": doc.metadata.updated_date,
            "raw_text_length": doc.raw_text_length,
            "total_tokens": doc.total_tokens,
            "num_chunks": len(doc.chunks),
            "chunks": [
                {
                    "chunk_id": chunk.chunk_id,
                    "content": chunk.content,
                    "tags": chunk.tags,
                    "live_lms_url": chunk.live_lms_url,
                    "token_count": chunk.token_count,
                    "chunk_index": chunk.chunk_index,
                    "total_chunks": chunk.total_chunks
                }
                for chunk in doc.chunks
            ]
        }
        
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(output_data, f, indent=2, ensure_ascii=False)
    
    def collect_files(self) -> List[Path]:
        """Collect all files to process."""
        files = []
        for ext in self.config.supported_extensions:
            files.extend(self.config.source_dir.rglob(f'*{ext}'))
        return sorted(files)
    
    def run(self):
        """Execute the full ingestion pipeline."""
        self.logger.logger.info("=" * 80)
        self.logger.logger.info("Luminate AI - Blackboard LMS Data Ingestion Pipeline")
        self.logger.logger.info("=" * 80)
        self.logger.logger.info(f"Source: {self.config.source_dir}")
        self.logger.logger.info(f"Output: {self.config.output_dir}")
        self.logger.logger.info(f"Course: {self.config.course_name} ({self.config.course_id})")
        
        # Parse manifest
        manifest_path = self.config.source_dir / 'imsmanifest.xml'
        if manifest_path.exists():
            self.logger.logger.info("Parsing manifest...")
            manifest_parser = ManifestParser(manifest_path, self.logger)
            structure = manifest_parser.parse()
        
        # Collect files
        self.logger.logger.info("Collecting files...")
        files = self.collect_files()
        self.stats['total_files'] = len(files)
        self.logger.logger.info(f"Found {len(files)} files to process")
        
        # Process files
        self.logger.logger.info("Processing files...")
        for file_path in tqdm(files, desc="Processing", unit="file"):
            try:
                doc = self.process_file(file_path)
                
                if doc:
                    self.save_document(doc)
                    
                    # Update statistics
                    self.stats['processed_files'] += 1
                    self.stats['total_chunks'] += len(doc.chunks)
                    self.stats['total_tokens'] += doc.total_tokens
                    self.stats['by_type'][doc.metadata.content_type] += 1
                    
                    module_stats = self.stats['by_module'][doc.metadata.module_name or 'Unknown']
                    module_stats['files'] += 1
                    module_stats['chunks'] += len(doc.chunks)
                    module_stats['tokens'] += doc.total_tokens
                else:
                    self.stats['skipped_files'] += 1
                    
            except Exception as e:
                self.stats['errors'] += 1
                self.logger.log_issue('ERROR', str(file_path), f'Processing failed: {e}')
        
        # Build sequential graph links
        self.logger.logger.info("Building graph relationships...")
        self.graph_builder.build_sequential_links()
        
        # Save graph
        graph_output = self.config.graph_dir / 'graph_links.json'
        self.graph_builder.save(graph_output)
        self.logger.logger.info(f"Saved {len(self.graph_builder.links)} graph links")
        
        # Save summary
        self.save_summary()
        
        # Save issues log
        self.logger.save_issues()
        
        self.logger.logger.info("=" * 80)
        self.logger.logger.info("Pipeline Complete!")
        self.logger.logger.info(f"Processed: {self.stats['processed_files']}/{self.stats['total_files']} files")
        self.logger.logger.info(f"Chunks: {self.stats['total_chunks']}")
        self.logger.logger.info(f"Tokens: {self.stats['total_tokens']:,}")
        self.logger.logger.info(f"Errors: {self.stats['errors']}")
        self.logger.logger.info("=" * 80)
    
    def save_summary(self):
        """Save ingestion summary."""
        summary = {
            "pipeline_info": {
                "run_date": datetime.now().isoformat(),
                "course_id": self.config.course_id,
                "course_name": self.config.course_name,
                "source_directory": str(self.config.source_dir),
                "output_directory": str(self.config.output_dir)
            },
            "statistics": {
                "total_files": self.stats['total_files'],
                "processed_files": self.stats['processed_files'],
                "skipped_files": self.stats['skipped_files'],
                "error_count": self.stats['errors'],
                "total_chunks": self.stats['total_chunks'],
                "total_tokens": self.stats['total_tokens'],
                "avg_tokens_per_file": (
                    self.stats['total_tokens'] / self.stats['processed_files']
                    if self.stats['processed_files'] > 0 else 0
                )
            },
            "by_file_type": dict(self.stats['by_type']),
            "by_module": {
                module: {
                    "files": stats['files'],
                    "chunks": stats['chunks'],
                    "tokens": stats['tokens']
                }
                for module, stats in self.stats['by_module'].items()
            }
        }
        
        summary_path = Path('ingest_summary.json')
        with open(summary_path, 'w', encoding='utf-8') as f:
            json.dump(summary, f, indent=2, ensure_ascii=False)
        
        self.logger.logger.info(f"Summary saved to {summary_path}")


# ============================================================================
# CLI ENTRY POINT
# ============================================================================

def main():
    """Main entry point for the script."""
    import argparse
    
    parser = argparse.ArgumentParser(
        description='Luminate AI - Blackboard LMS Data Ingestion Pipeline',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # Process with default settings
  python ingest_clean_luminate.py
  
  # Specify custom source directory
  python ingest_clean_luminate.py --source /path/to/export
  
  # Specify custom course ID
  python ingest_clean_luminate.py --course-id _12345_1
        """
    )
    
    parser.add_argument(
        '--source',
        type=Path,
        default=Path('extracted/ExportFile_COMP237_INP'),
        help='Source directory containing Blackboard export'
    )
    
    parser.add_argument(
        '--output',
        type=Path,
        default=Path('cleaned'),
        help='Output directory for cleaned JSON files'
    )
    
    parser.add_argument(
        '--course-id',
        type=str,
        default='_29430_1',
        help='Blackboard course ID'
    )
    
    parser.add_argument(
        '--course-name',
        type=str,
        default='Luminate',
        help='Course name'
    )
    
    args = parser.parse_args()
    
    # Create configuration
    config = Config(
        source_dir=args.source,
        output_dir=args.output,
        course_id=args.course_id,
        course_name=args.course_name
    )
    
    # Validate source directory
    if not config.source_dir.exists():
        print(f"Error: Source directory does not exist: {config.source_dir}")
        exit(1)
    
    # Run pipeline
    pipeline = IngestionPipeline(config)
    pipeline.run()


if __name__ == '__main__':
    main()
